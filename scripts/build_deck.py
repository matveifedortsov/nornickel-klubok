"""Сборка презентации из HTML-исходника в PNG + PDF + PPTX.

Дизайн живёт в docs/presentation.html (дизайн-система с кастомными SVG). Этот
скрипт рендерит каждый слайд headless-браузером в высоком разрешении и собирает:
  - runtime/slides/s01..NN.png  — по слайду
  - docs/presentation.pdf       — по картинке на страницу (16:9 landscape)
  - docs/presentation.pptx      — full-bleed картинка на слайд + заметки докладчика

Так дизайн переносится 1:1 (текст — растр, зато пиксель-в-пиксель как в HTML).
Требует: playwright (+ chromium или Edge), Pillow, python-pptx.

Запуск:  python scripts/build_deck.py
"""
from __future__ import annotations

from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
HTML = ROOT / "docs" / "presentation.html"
SLIDES = ROOT / "runtime" / "slides"
PDF = ROOT / "docs" / "presentation.pdf"
PPTX = ROOT / "docs" / "presentation.pptx"
SCALE = 2  # device_scale_factor: 1280x720 -> 2560x1440

# Заметки докладчика по слайдам (Presenter View). Порядок = порядок секций.
NOTES = [
    # 1 Титул
    "Здравствуйте! Команда трека «Научный клубок». Наш тезис: проблема R&D не в "
    "отсутствии моделей, а в отсутствии структуры знаний. Мы строим граф знаний + "
    "гибридный поиск, и решение LLM-агностично: сегодня YandexGPT, в проде тот же "
    "код с MetalGPT-1 on-prem. План: 30 сек проблема, минута архитектура, дальше демо.",
    # 2 Проблема
    "Пять болей — дословно из вашего ТЗ. Знания рассеяны: команды дублируют "
    "литобзоры, нет междисциплинарных связей, многопараметрический поиск невозможен, "
    "противоречия без единой базы, численные фильтры теряются. Ответ на один "
    "техвопрос сегодня = дни ручного сбора. Мы сводим к секундам — с провенансом.",
    # 3 Решение
    "Связываем публикации, эксперименты, материалы, процессы, оборудование, "
    "экспертов в единый граф. Вектор находит смысл, граф даёт структуру и "
    "доказательную базу. Ключевое отличие — не «RAG поверх PDF»: каждый ответ с "
    "цитатой, уровнем достоверности и датой. Фильтры распознаются прямо из вопроса.",
    # 4 Онтология
    "Домен-специфичная онтология: 12 типов сущностей, 15 типов связей. LLM извлекает "
    "триплеты строго по схеме в JSON, каждое ребро несёт evidence-цитату. Сейчас в "
    "графе 2111 узлов и 1436 рёбер — из отраслевых журналов и open-access статей. "
    "На схеме — типичный фрагмент: эксперимент применяет процесс к материалу и т.д.",
    # 5 Архитектура
    "Главный козырь — LLM-агностичность. Один пайплайн: парсинг → извлечение LLM в "
    "JSON → Neo4j + Qdrant. Запрос: вектор → seed → обход графа 1–4 хопа с фильтрами "
    "→ реранк → генерация. Бэкенд меняется строкой в .env: YandexGPT сегодня, "
    "MetalGPT-1 on-prem в проде — данные не покидают периметр. Эмбеддинги локальные, без GPU.",
    # 6 Гибридный поиск
    "Пять ступеней. Важно: структурные фильтры — числа, гео, годы — применяются "
    "ВНУТРИ обхода графа, до генерации. Ретривал медиана 447 мс при бюджете ТЗ "
    "3–5 секунд: запас в разы. Пример многопараметрического вопроса про "
    "обессоливание с числовыми диапазонами — распознаётся и применяется к подграфу.",
    # 7 Метрики
    "Всё измерено на реальном стенде, не на глаз. Held-out gold не пересекается с "
    "few-shot. F1: сущности 0,66, связи 0,62, атрибуты 0,67. Честно про слабые "
    "места: gold пока 10 примеров, кросс-язычное сопоставление меток — известный "
    "пробел. Мягкая деградация проверена: ни один сбой не роняет ответ целиком.",
    # 8 Покрытие ТЗ
    "Все блоки ТЗ закрыты рабочим кодом с тестами: онтология, многопараметрический "
    "поиск, верификация, аналитика (пробелы, литобзор, сравнение, эксперты), RBAC + "
    "аудит, правка графа, экспорт, уведомления, дашборд. Нефункциональные тоже: "
    "ретривал в бюджете, без GPU, 68 тестов, RU/EN.",
    # 9 Развитие
    "Готово к пилоту в НИИ уже сейчас. Дорожная карта: on-prem MetalGPT-1 (клиент "
    "написан), полный корпус, каталог экспериментов. Почему работает в проде: "
    "LLM-агностичность, периметр безопасности, провенанс by design, мягкая "
    "деградация. docker compose up — и жюри видит рабочее демо за минуту. Спасибо!",
]

# CSS для экспорта: расплющиваем секции в чистые прямоугольники 1280x720
# (без border-radius/тени/полей — иначе по углам просвечивает фон body).
FLATTEN_CSS = """
  html, body { padding: 0 !important; margin: 0 !important; background: #062B4A !important; }
  section { margin: 0 !important; border-radius: 0 !important; box-shadow: none !important;
            transform: none !important; }
  #progress-bar, #slide-counter, #nav-dots, .arrow { display: none !important; }
"""


def render_slides() -> list[Path]:
    from playwright.sync_api import sync_playwright
    SLIDES.mkdir(parents=True, exist_ok=True)
    for old in SLIDES.glob("s*.png"):
        old.unlink()
    out: list[Path] = []
    with sync_playwright() as p:
        try:
            b = p.chromium.launch()
        except Exception:                                   # noqa: BLE001 — нет chromium
            b = p.chromium.launch(channel="msedge")
        pg = b.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=SCALE)
        pg.goto(HTML.resolve().as_uri())
        pg.add_style_tag(content=FLATTEN_CSS)
        pg.wait_for_timeout(1200)
        secs = pg.query_selector_all("section")
        for i, s in enumerate(secs, 1):
            s.scroll_into_view_if_needed()
            pg.wait_for_timeout(250)
            path = SLIDES / f"s{i:02d}.png"
            s.screenshot(path=str(path))
            out.append(path)
        b.close()
    print(f"отрендерено слайдов: {len(out)}")
    return out


def build_pdf(pngs: list[Path]) -> None:
    from PIL import Image
    imgs = [Image.open(p).convert("RGB") for p in pngs]
    imgs[0].save(PDF, save_all=True, append_images=imgs[1:], resolution=150.0)
    print(f"PDF -> {PDF} ({len(imgs)} стр.)")


def build_pptx(pngs: list[Path]) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i, png in enumerate(pngs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        if i < len(NOTES):
            slide.notes_slide.notes_text_frame.text = NOTES[i]
    prs.save(PPTX)
    print(f"PPTX -> {PPTX} ({len(prs.slides._sldIdLst)} слайдов)")


def main() -> None:
    pngs = render_slides()
    if not pngs:
        raise SystemExit("не удалось отрендерить слайды")
    build_pdf(pngs)
    build_pptx(pngs)
    print("готово.")


if __name__ == "__main__":
    main()
