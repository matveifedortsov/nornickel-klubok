"""PPTX -> Document(chunks). Тот же контракт, что и pdf_parser.parse_pdf.

Доклады/презентации в корпусе несут текст по слайдам, включая таблицы и
заметки докладчика (notes) — последние часто содержат больше содержательного
текста, чем сами слайды. Тяжёлая библиотека (python-pptx) импортируется
лениво, как и в pdf_parser.py.
"""
from __future__ import annotations

from pathlib import Path

from klubok.ontology import Document, Chunk
from klubok.parsing.pdf_parser import chunk_text, _doc_id


def _shape_text(shape) -> list[str]:
    out: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs) or para.text
            if text.strip():
                out.append(text.strip())
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append("\t".join(cells))
    return out


def parse_pptx(path: str | Path, max_chars: int = 1200) -> Document:
    """Распарсить один PPTX в Document. Требует python-pptx (ленивый импорт)."""
    from pptx import Presentation

    path = Path(path)
    doc_id = _doc_id(path)
    document = Document(doc_id=doc_id, title=path.stem, source_path=str(path))

    prs = Presentation(str(path))
    chunks: list[Chunk] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_text(shape))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append(slide.notes_slide.notes_text_frame.text.strip())
        text = "\n\n".join(parts)
        chunks.extend(chunk_text(text, doc_id, page=slide_no, max_chars=max_chars))

    document.chunks = chunks
    return document


def parse_dir(folder: str | Path, pattern: str = "*.pptx") -> list[Document]:
    folder = Path(folder)
    return [parse_pptx(p) for p in sorted(folder.glob(pattern))]
